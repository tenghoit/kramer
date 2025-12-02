from pypdf import PdfReader, errors
from pptx import Presentation
from docx import Document
from pathlib import Path
import re
import sys
import src.setup_logging

def pdf_to_text(file_path: Path) -> str:
    try:
        reader = PdfReader(file_path)
        text = "\n".join(page.extract_text() or "" for page in reader.pages) # or "" b/c reader return None if image
        text = re.sub(r"[\n]{2,}", "\n", text) # Replace multiple newlines with a single newline
        text = re.sub(r"[\s]+", " ", text) # Replace multiple spaces with a single space
        return text 
    except errors.PyPdfError as e:
        raise ValueError(f"Failed to read PDF: {file_path}") from e


def docx_to_text(file_path: Path) -> str:
    try:
        document = Document(str(file_path))
        return "\n".join(paragraph.text for paragraph in document.paragraphs)
    except Exception as e:
        raise ValueError(f"Failed to read DOCX: {file_path}") from e


def extract_text(file_path: str | Path):
    """
    Extract text from a PDF or DOCX file.
    """
    
    file_path = Path(file_path.strip()) if isinstance(file_path, str) else file_path

    if not file_path.is_file():
        raise FileNotFoundError(f"Invalid file path: {file_path}")

    suffix = file_path.suffix.lower()
    if suffix == ".txt":
        return file_path.read_text()        
    elif suffix == ".pdf":
        return pdf_to_text(file_path)
    elif suffix == ".docx":
        return docx_to_text(file_path)
    elif suffix == ".pptx":
        return pptx_to_texts(file_path)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")
    
def pptx_to_texts(file_path) -> list[str]: 
    prs = Presentation(file_path)
    slide_texts = []

    for slide in prs.slides:
        slide_text = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # text_runs.append(run.text)
                    slide_text += run.text
        slide_texts.append(slide_text)

    return slide_texts
    

def main():
    
    if len(sys.argv) != 2:
        print(f"Usage: python3 {sys.argv[0]}.py <file_path>")
        sys.exit(1)
    
    # print(extract_text(sys.argv[1]))

    for text in extract_text(sys.argv[1]):
        print(f"{text}\n<page-break>")


if __name__ == "__main__": 
    main()